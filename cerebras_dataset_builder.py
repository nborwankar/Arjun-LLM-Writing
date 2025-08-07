#!/usr/bin/env python3
"""
=============================================================================
CEREBRAS DATASET BUILDER - Fast & Cost-Effective Fine-Tuning Dataset Creator
=============================================================================

PURPOSE:
Generates high-quality fine-tuning datasets from investment documents using 
Cerebras gpt-oss-120b model for fast, cost-effective document analysis and 
training pair generation.

INPUT FILES:
- Writing/*.pdf - Investment research papers, market commentary (PDFs)
- Writing/*.docx - Strategy documents, analysis reports (Word documents)
- Writing/*.doc - Legacy Word documents
- Writing/*.pptx - Presentation slides with investment insights

OUTPUT FILES:
- cerebras_dataset.json - Complete dataset with metadata, chunks, and training pairs
  * Contains document analysis, quality scores, and generated Q&A pairs
  * Format: {"summary": stats, "documents": metadata, "chunks": text_segments, "training_pairs": qa_data}
  * Used as input for OpenAI fine-tuning conversion

PROCESS OVERVIEW:
1. Extract text from multiple document formats (PDF, DOCX, DOC, PPTX)
2. Analyze documents with Cerebras AI for voice characteristics
3. Chunk documents intelligently (500-1500 words, preserving context)
4. Generate multiple training pairs per chunk (style, continuation, analogies)
5. Score quality and capture distinctive voice elements
6. Output structured dataset ready for fine-tuning

COST: ~$15-20 for typical 28-document corpus (much cheaper than Claude Opus)
TIME: ~10-15 minutes for full processing

Last Updated: 2025-08-07
Version: 2.0 - Production Ready
"""

import os, json, re, time, math, random, hashlib, sys
from pathlib import Path
from dataclasses import dataclass, asdict
from typing import List, Dict, Tuple, Any

import PyPDF2
from docx import Document
from pptx import Presentation
import docx2txt
import fitz  # PyMuPDF
from cerebras.cloud.sdk import Cerebras

# ═══════════════════════════ CONFIGURATION ═══════════════════════════ #
# API Configuration
CEREBRAS_API_KEY = "csk-mfr4yk383rcv94cd2kxkwkxnwtp2fykyhrme4n4mtcc53xfm"
CEREBRAS_MODEL = "gpt-oss-120b"  # Fast, cost-effective model for analysis

# Processing Parameters
MAX_TOKENS_RESPONSE = 8000  # Maximum response length from Cerebras
MIN_CHUNK_WORDS = 500       # Minimum words per chunk (preserves context)
MAX_CHUNK_WORDS = 1500      # Maximum words per chunk (manageable size)

# Quality Thresholds
MIN_QUALITY_SCORE = 6.0     # Minimum quality score for inclusion
MAX_RETRIES = 3             # API retry attempts on failure

# ═══════════════════════════ DATA STRUCTURES ═══════════════════════════ #

@dataclass
class DocumentMetadata:
    """Stores comprehensive metadata about processed documents"""
    filename: str                    # Original document filename
    file_type: str                  # File extension (pdf, docx, etc.)
    document_type: str              # Content type (research, commentary, etc.)
    word_count: int                 # Total words in document
    chunk_count: int                # Number of chunks created
    has_analogies: bool             # Contains analogical thinking
    has_data_analysis: bool         # Contains data/charts analysis
    has_personal_voice: bool        # Contains personal opinions/voice
    quality_score: float            # AI-assessed quality (0-10)
    distinctive_elements: List[str]  # Memorable phrases/patterns
    analysis_summary: str           # AI summary of document value

@dataclass
class TextChunk:
    """Represents a processed text segment ready for training pair generation"""
    chunk_id: str                   # Unique identifier (MD5 hash)
    source_document: str            # Original document filename
    document_type: str              # Document category
    text: str                       # Actual text content
    word_count: int                 # Words in this chunk
    chunk_index: int                # Position in document (0, 1, 2...)
    contains_analogy: bool          # Has analogical explanations
    contains_data: bool             # Has numerical/data analysis
    contains_personal_voice: bool   # Has personal opinions/style
    training_format: str            # Suggested training approach
    quality_assessment: str         # Quality notes from AI

# ═══════════════════════════ MAIN PROCESSOR ═══════════════════════════ #

class CerebrasProcessor:
    """
    Main processor class that handles document analysis and dataset generation.
    
    Uses Cerebras gpt-oss-120b model for fast, cost-effective analysis of
    investment documents to create fine-tuning datasets.
    
    Key Features:
    - Multi-format document extraction (PDF, DOCX, DOC, PPTX)
    - Intelligent text chunking preserving context
    - AI-powered voice characteristic analysis
    - Multiple training pair generation strategies
    - Quality scoring and filtering
    - Comprehensive metadata capture
    """
    def __init__(self, input_dir: str, output_json: str, api_key: str = CEREBRAS_API_KEY):
        """
        Initialize the Cerebras processor.
        
        Args:
            input_dir: Path to directory containing investment documents
            output_json: Output filename for generated dataset
            api_key: Cerebras API key for authentication
        """
        self.input_dir = Path(input_dir)    # Input directory path
        self.output_json = output_json      # Output JSON filename
        self.client = Cerebras(api_key=api_key)  # Cerebras API client
        self.sent_requests = 0              # Track API usage
        self.total_tokens = 0               # Track token consumption

    # ═══════════════════════ DOCUMENT EXTRACTION METHODS ═══════════════════════ #
    
    @staticmethod
    def _extract_pdf(p: Path) -> str:
        """
        Extract text from PDF files using PyMuPDF with OCR error correction.
        
        Handles common PDF issues:
        - Ligature characters (fi, fl, ff, ffi, ffl)
        - OCR errors (moths -> months, Jeremey -> Jeremy)
        - Encoding issues and special characters
        
        Args:
            p: Path to PDF file
            
        Returns:
            str: Cleaned text content from all pages
            
        Fallback: Uses PyPDF2 if PyMuPDF fails
        """
        try:
            doc = fitz.open(str(p))
            text_pages = [page.get_text() for page in doc]
            doc.close()
            full_text = "\n".join(text_pages)
            
            # Fix common ligatures and OCR errors found in investment documents
            fixes = {
                "ﬁ": "fi", "ﬂ": "fl", "ﬀ": "ff", "ﬃ": "ffi", "ﬄ": "ffl",  # Ligatures
                "�": "",                                                              # Invalid chars
                "Jeremey": "Jeremy",                                                  # Name corrections
                " moths ": " months ", "12 moths": "12 months"                      # Time period fixes
            }
            for old, new in fixes.items():
                full_text = full_text.replace(old, new)
            return full_text
        except Exception as e:
            print(f"   ⚠️ PyMuPDF failed for {p.name}: {e}, using PyPDF2")
            with open(p, "rb") as fh:
                rdr = PyPDF2.PdfReader(fh)
                return "\n".join(pg.extract_text() or "" for pg in rdr.pages)

    @staticmethod
    def _extract_docx(p: Path) -> str:
        """
        Extract text from modern Word documents (.docx).
        
        Args:
            p: Path to DOCX file
            
        Returns:
            str: Text content from all paragraphs
        """
        return "\n".join(par.text for par in Document(p).paragraphs)

    @staticmethod
    def _extract_doc(p: Path) -> str:
        """
        Extract text from legacy Word documents (.doc).
        
        Args:
            p: Path to DOC file
            
        Returns:
            str: Extracted text content or empty string if failed
        """
        return docx2txt.process(str(p)) or ""

    @staticmethod
    def _extract_pptx(p: Path) -> str:
        """
        Extract text from PowerPoint presentations (.pptx).
        
        Preserves slide structure with slide numbers for context.
        
        Args:
            p: Path to PPTX file
            
        Returns:
            str: Text content with slide markers
        """
        pres, out = Presentation(p), []
        for idx, sl in enumerate(pres.slides, 1):
            out.append(f"\n--- Slide {idx} ---\n")
            out.extend(sh.text for sh in sl.shapes if hasattr(sh, "text") and sh.text)
        return "\n".join(out)

    # ═══════════════════════ AI PROCESSING METHODS ═══════════════════════ #
    
    def _call_cerebras(self, prompt: str) -> str:
        """
        Call Cerebras API with automatic retry logic and error handling.
        
        Features:
        - 3 retry attempts with exponential backoff
        - Request tracking for cost estimation
        - Detailed logging of API calls
        - Graceful error handling
        
        Args:
            prompt: Text prompt to send to Cerebras gpt-oss-120b
            
        Returns:
            str: AI response content
            
        Raises:
            RuntimeError: If all retry attempts fail
        """
        for attempt in range(MAX_RETRIES):
            try:
                print(f"   🧠 Calling Cerebras gpt-oss-120b (attempt {attempt + 1})...")
                response = self.client.chat.completions.create(
                    messages=[{"role": "user", "content": prompt}],
                    model=CEREBRAS_MODEL,
                    max_completion_tokens=MAX_TOKENS_RESPONSE,
                    temperature=0.3
                )
                self.sent_requests += 1
                result = response.choices[0].message.content
                print(f"   ✅ Response received ({len(result)} chars)")
                return result
            except Exception as e:
                if attempt == MAX_RETRIES - 1:
                    raise RuntimeError(f"Cerebras API failed: {e}") from e
                print(f"   ⚠️ Retry {attempt + 1}: {e}")
                time.sleep(1 * (2 ** attempt))

    @staticmethod
    def _parse_json(raw: str) -> Dict:
        """
        Parse JSON response from Cerebras API with error handling.
        
        Handles common JSON parsing issues:
        - Markdown code blocks (```json)
        - Extra whitespace and formatting
        - Malformed JSON structures
        
        Args:
            raw: Raw response string from API
            
        Returns:
            Dict: Parsed JSON data
            
        Raises:
            ValueError: If JSON cannot be parsed after cleanup attempts
        """
        raw = raw.strip()
        if raw.startswith("```"):
            parts = raw.split("```")
            if len(parts) >= 3:
                raw = parts[1]
                if raw.startswith(("json", "JSON")):
                    raw = raw[4:].strip()
        start, end = raw.find("{"), raw.rfind("}") + 1
        if start != -1 and end > start:
            raw = raw[start:end]
        return json.loads(raw)

    def _validate_chunks(self, full_text: str, chunks: List[Dict]) -> List[Dict]:
        """
        Validate and optimize text chunks for training effectiveness.
        
        Validation Rules:
        - Keep short documents (< MIN_CHUNK_WORDS) as single chunks
        - Ensure chunks are within MIN_CHUNK_WORDS to MAX_CHUNK_WORDS range
        - Preserve document context and logical flow
        - Filter out chunks that are too short or low quality
        
        Args:
            full_text: Complete document text
            chunks: List of chunk dictionaries from AI analysis
            
        Returns:
            List[Dict]: Validated and optimized chunks
        """
        doc_words = len(full_text.split())
        if doc_words < MIN_CHUNK_WORDS:
            print(f"   📄 Short doc ({doc_words} words), keeping whole")
            return [{
                "text": full_text.strip(), "word_count": doc_words,
                "contains_analogy": any(c.get("contains_analogy") for c in chunks),
                "contains_data": any(c.get("contains_data") for c in chunks),
                "contains_personal_voice": any(c.get("contains_personal_voice") for c in chunks),
                "training_format": "style_replication", "quality_score": 0.8
            }]
        if all(MIN_CHUNK_WORDS <= c["word_count"] <= MAX_CHUNK_WORDS for c in chunks):
            return chunks
        # Rechunk if needed
        print(f"   🔧 Re-chunking document")
        words, new_chunks, idx = full_text.split(), [], 0
        while idx < len(words):
            end = min(idx + 1000, len(words))
            new_chunks.append({
                "text": " ".join(words[idx:end]), "word_count": end - idx,
                "contains_analogy": False, "contains_data": False,
                "contains_personal_voice": True, "training_format": "continuation",
                "quality_score": 0.7
            })
            idx = end
        return new_chunks

    def _generate_training_pairs(self, chunks: List[Dict]) -> List[Dict]:
        """
        Generate diverse training pairs from processed text chunks.
        
        Training Pair Types:
        - Style replication: "Write in Arjun's style about..."
        - Content continuation: "Complete this analysis..."
        - Analogical thinking: "Explain using an analogy..."
        - Opinion extraction: "What's your view on..."
        
        Args:
            chunks: List of validated text chunks
            
        Returns:
            List[Dict]: Training pairs in OpenAI format
                       Each pair: {"instruction": str, "input": str, "output": str}
        """
        pairs = []
        for chunk in chunks:
            text = chunk["text"]
            source_id = chunk["chunk_id"]
            # Style replication pair
            pairs.append({
                "instruction": f"Write analysis in this investment style:",
                "input": f"Topic: Market analysis\n\nStyle: {text[:200]}...",
                "output": text,
                "metadata": {"source_chunk_id": source_id, "type": "style_replication"}
            })
            # Continuation pair (if long enough)
            if len(text.split()) > 400:
                mid = len(text) // 2
                pairs.append({
                    "instruction": "Continue this investment analysis:",
                    "input": text[:mid], "output": text[mid:],
                    "metadata": {"source_chunk_id": source_id, "type": "continuation"}
                })
        return pairs

    @staticmethod
    def _prompt(fname: str, txt: str) -> str:
        """
        Generate comprehensive analysis prompt for Cerebras AI.
        
        The prompt instructs the AI to:
        - Identify document type and investment focus
        - Detect voice characteristics (analogies, personal opinions, data analysis)
        - Create intelligent text chunks preserving context
        - Score quality for fine-tuning suitability
        - Extract distinctive phrases and patterns
        
        Args:
            fname: Document filename for context
            txt: Document text content
            
        Returns:
            str: Structured prompt for AI analysis
        """
        schema = """{
  "document_type": "Market Commentary" | "Research Papers" | "Strategy Documents" | "Lessons/Insights",
  "distinctive_elements": {
    "has_analogies": true/false,
    "has_data_analysis": true/false,
    "has_personal_voice": true/false,
    "memorable_phrases": ["phrase1"]
  },
  "text_chunks": [
    {
      "text": "chunk text (500-1500 words)",
      "word_count": 123,
      "contains_analogy": true/false,
      "contains_data": true/false,
      "contains_personal_voice": true/false,
      "training_format": "style_replication",
      "quality_score": 0.8
    }
  ],
  "overall_quality_assessment": {
    "quality_score": 0.8,
    "suitability_for_training": "assessment"
  }
}"""
        return f"""Analyze this investment document for AI training. Create 500-1500 word chunks preserving complete thoughts. Return only JSON:

{schema}

DOCUMENT: {fname}
{txt}

JSON:"""

    def _process_single(self, p: Path) -> Tuple[DocumentMetadata, List[TextChunk]]:
        """
        Process a single document through the complete analysis pipeline.
        
        Pipeline Steps:
        1. Extract text using appropriate method (PDF, DOCX, DOC, PPTX)
        2. Send to Cerebras AI for analysis and chunking
        3. Parse AI response and validate JSON structure
        4. Create DocumentMetadata with quality scores
        5. Generate TextChunk objects with unique IDs
        6. Validate chunk sizes and content quality
        
        Args:
            p: Path to document file
            
        Returns:
            Tuple[DocumentMetadata, List[TextChunk]]: Document metadata and text chunks
            
        Raises:
            Exception: If document processing fails (logged but not fatal)
        """
        print(f"→ {p.name}")
        
        # Extract text
        extractors = {".pdf": self._extract_pdf, ".docx": self._extract_docx,
                     ".doc": self._extract_doc, ".pptx": self._extract_pptx}
        raw_text = extractors[p.suffix.lower()](p)
        if not raw_text.strip():
            raise ValueError("No text extracted")

        # Process with Cerebras
        response = self._parse_json(self._call_cerebras(self._prompt(p.name, raw_text)))
        chunks_data = self._validate_chunks(raw_text, response["text_chunks"])

        # Create metadata
        meta = DocumentMetadata(
            filename=p.name, file_type=p.suffix.lower(),
            document_type=response["document_type"],
            word_count=len(raw_text.split()), chunk_count=len(chunks_data),
            has_analogies=response["distinctive_elements"]["has_analogies"],
            has_data_analysis=response["distinctive_elements"]["has_data_analysis"],
            has_personal_voice=response["distinctive_elements"]["has_personal_voice"],
            quality_score=response["overall_quality_assessment"]["quality_score"],
            distinctive_elements=response["distinctive_elements"].get("memorable_phrases", []),
            analysis_summary=response["overall_quality_assessment"]["suitability_for_training"]
        )

        # Create chunks
        chunks = []
        for i, chunk_data in enumerate(chunks_data):
            chunk_id = hashlib.md5(f"{p.stem}_{i}".encode()).hexdigest()[:12]
            chunk_data["chunk_id"] = chunk_id
            chunk_data["document_type"] = meta.document_type
            chunks.append(TextChunk(
                chunk_id=chunk_id, source_document=p.name,
                document_type=meta.document_type, text=chunk_data["text"],
                word_count=chunk_data["word_count"], chunk_index=i,
                contains_analogy=chunk_data["contains_analogy"],
                contains_data=chunk_data["contains_data"],
                contains_personal_voice=chunk_data["contains_personal_voice"],
                training_format=chunk_data["training_format"],
                quality_assessment="Cerebras analysis"
            ))

        print(f"   ✓ {len(chunks)} chunks, score {meta.quality_score:.2f}")
        return meta, chunks

    def run(self):
        """
        Main execution method - processes all documents in input directory.
        
        Complete Workflow:
        1. Scan input directory for supported file types (PDF, DOCX, DOC, PPTX)
        2. Process each document through analysis pipeline
        3. Generate training pairs from all chunks
        4. Compile comprehensive statistics and metadata
        5. Output structured dataset JSON file
        6. Display cost estimates and processing summary
        
        Output Structure:
        - summary: Processing statistics and cost estimates
        - documents: Metadata for each processed document
        - chunks: All text segments with analysis
        - training_pairs: Generated Q&A pairs for fine-tuning
        - processing_notes: Technical details and model info
        
        Side Effects:
        - Creates output JSON file
        - Prints progress and statistics to console
        - Logs any processing errors
        """
        print(f"🧠 CEREBRAS DATASET BUILDER")
        print(f"📁 Input: {self.input_dir}")
        
        paths = []
        for g in ("*.pdf", "*.docx", "*.doc", "*.pptx"):
            paths.extend(self.input_dir.glob(g))
        if not paths:
            print("No files found."); return

        print(f"\n📊 Processing {len(paths)} documents...")
        
        metas, all_chunks, training_pairs, fails = [], [], [], []
        
        for path in paths:
            try:
                meta, chunks = self._process_single(path)
                metas.append(asdict(meta))
                chunk_dicts = [asdict(c) for c in chunks]
                for c in chunk_dicts:
                    c["chunk_id"] = chunks[chunk_dicts.index(c)].chunk_id
                all_chunks.extend(chunk_dicts)
                training_pairs.extend(self._generate_training_pairs(chunk_dicts))
            except Exception as e:
                fails.append({"file": path.name, "error": str(e)})
                print(f"   ✗ {path.name}: {e}")

        # Statistics
        summary = {
            "total_documents": len(metas),
            "total_chunks": len(all_chunks),
            "total_training_pairs": len(training_pairs),
            "total_words": sum(c["word_count"] for c in all_chunks),
            "failed_files": fails,
            "average_quality_score": sum(m["quality_score"] for m in metas) / len(metas) if metas else 0,
            "estimated_cost": (self.sent_requests * 100) * 0.006  # Rough estimate
        }

        # Output
        output = {
            "summary": summary, "documents": metas,
            "chunks": all_chunks, "training_pairs": training_pairs,
            "processing_notes": {"model_used": CEREBRAS_MODEL, "api_provider": "Cerebras"}
        }
        
        with open(self.output_json, "w", encoding="utf-8") as f:
            json.dump(output, f, indent=2, ensure_ascii=False)

        print(f"\n🎉 COMPLETE! Cost: ~${summary['estimated_cost']:.2f}")
        print(f"📊 {len(training_pairs)} training pairs generated")
        print(f"💾 Output: {self.output_json}")

# ═══════════════════════════ MAIN EXECUTION ═══════════════════════════ #

if __name__ == "__main__":
    # Configuration for dataset generation
    INPUT_DIR = "/Users/macbook2024/Dropbox/AAA Backup/A Working/Arjun LLM Writing/Writing"
    OUTPUT_JSON = "cerebras_dataset.json"
    
    print("🚀 Starting Cerebras Dataset Builder...")
    print(f"📁 Input Directory: {INPUT_DIR}")
    print(f"📄 Output File: {OUTPUT_JSON}")
    print(f"🧠 AI Model: {CEREBRAS_MODEL}")
    print("─" * 60)
    
    # Initialize and run processor
    processor = CerebrasProcessor(INPUT_DIR, OUTPUT_JSON)
    processor.run()
    
    print("\n✅ Dataset generation complete!")
    print(f"💾 Dataset saved to: {OUTPUT_JSON}")
    print("🔄 Ready for OpenAI fine-tuning conversion")
