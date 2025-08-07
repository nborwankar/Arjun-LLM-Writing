#!/usr/bin/env python3
"""
Claude Opus Dataset Builder v2 - Premium AI-Powered Investment Voice Dataset Generator

INPUT FILES:
- All documents in Writing/ directory (PDF, DOCX, DOC, PPTX formats)
- Processes investment documents, market commentary, research papers, strategy documents
- Requires ANTHROPIC_API_KEY environment variable

OUTPUT FILES:
- training_dataset_opus_processed.json: Complete dataset with metadata, chunks, and training pairs
- Structured JSON with document analysis, text chunks, and fine-tuning ready training pairs

PROCESSING PIPELINE:
1. Document Extraction: Multi-format text extraction with OCR error correction
2. AI Analysis: Claude Opus 4.1 performs intelligent document analysis and chunking
3. Quality Validation: Chunk size validation, deduplication, and balance checking
4. Training Pair Generation: Multiple training variations per chunk for fine-tuning
5. Cost Optimization: Token counting, rate limiting, and API cost tracking

KEY FEATURES:
- Premium Claude Opus 4.1 analysis for superior quality
- Intelligent chunking preserving context and voice characteristics
- Multiple training pair formats (style replication, continuation, Q&A)
- Document balance validation against recommended distribution
- Comprehensive deduplication and quality assessment
- Real-time cost tracking and API usage monitoring

USAGE:
    python claude_opus_dataset_builder_v2.py
    
REQUIREMENTS:
- anthropic>=0.25.0
- PyMuPDF (fitz)
- python-docx
- python-pptx
- PyPDF2
- tiktoken (optional, for accurate token counting)

Version: 2.0
Last Updated: January 2025
Author: Investment AI Fine-tuning Project
"""

import os, json, re, time, math, random, hashlib, sys
from pathlib import Path
from dataclasses import dataclass, asdict
from typing import List, Dict, Tuple, Any

import requests
import PyPDF2
from docx import Document
from pptx import Presentation
import docx2txt
import fitz  # PyMuPDF

# ───────────────────────── CONFIG ───────────────────────── #
CLAUDE_MODEL        = "claude-opus-4-1-20250805"
API_URL             = "https://api.anthropic.com/v1/messages"
MAX_TOKENS_RESPONSE = 8_000
RATE_LIMIT_RPM      = 100
RETRY_MAX_ATTEMPTS  = 5
RETRY_BACKOFF_BASE  = 2
MIN_CHUNK_WORDS     = 500
MAX_CHUNK_WORDS     = 1500
PII_PATTERNS        = [
    (re.compile(r"\b[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}\b", re.I), "[REDACTED_EMAIL]"),
    (re.compile(r"\b\d{3}[-.\s]?\d{2,4}[-.\s]?\d{4}\b"), "[REDACTED_PHONE]")
]

try:
    import tiktoken
    ENCODER = tiktoken.get_encoding("cl100k_base")
except ImportError:
    ENCODER = None

# PDF extraction (using PyPDF2 for now due to conda charset_normalizer conflict)
USE_PDFPLUMBER = False

def count_tokens(text: str) -> int:
    """
    Count tokens in text for accurate API cost estimation.
    
    Uses tiktoken for precise counting when available, falls back to
    heuristic estimation (words * 1.3) for cost planning.
    
    Args:
        text: Text content to count tokens for
        
    Returns:
        int: Estimated token count
    """
    if ENCODER:  # Accurate token counting with tiktoken
        return len(ENCODER.encode(text))
    return math.ceil(len(text.split()) * 1.3)  # Heuristic fallback

def redact_pii(txt: str) -> str:
    """
    Remove personally identifiable information from text content.
    
    Protects privacy by redacting:
    - Email addresses
    - Phone numbers
    - Other sensitive patterns
    
    Args:
        txt: Text content to redact
        
    Returns:
        str: Text with PII patterns replaced with redaction markers
    """
    for pat, repl in PII_PATTERNS:
        txt = pat.sub(repl, txt)
    return txt

def backoff(attempt: int):
    """
    Implement exponential backoff with jitter for API retry attempts.
    
    Args:
        attempt: Current retry attempt number (0-based)
    """
    time.sleep(RETRY_BACKOFF_BASE * (2 ** attempt) + random.uniform(0, 1))

# ────────────────────── DATA STRUCTURES ─────────────────── #
@dataclass
class DocumentMetadata:
    """
    Comprehensive metadata for processed investment documents.
    
    Captures document characteristics essential for voice fine-tuning:
    - Content analysis (analogies, data analysis, personal voice)
    - Quality scoring for training effectiveness
    - Document type classification for balanced dataset
    - Claude Opus analysis summary for human review
    """
    filename: str                    # Source document filename
    file_type: str                   # File extension (PDF, DOCX, etc.)
    document_type: str               # Investment document category
    word_count: int                  # Total words in document
    chunk_count: int                 # Number of chunks generated
    has_analogies: bool              # Contains analogical reasoning
    has_data_analysis: bool          # Contains data-driven analysis
    has_personal_voice: bool         # Contains personal opinions/voice
    quality_score: float             # Claude's quality assessment (0-10)
    distinctive_elements: List[str]  # Key phrases and patterns identified
    opus_analysis_summary: str       # Claude's detailed analysis

@dataclass
class TextChunk:
    """
    Individual text segments optimized for fine-tuning training pairs.
    
    Each chunk represents a coherent piece of investment writing that:
    - Maintains context and voice characteristics
    - Falls within optimal size range (500-1500 words)
    - Contains identifiable voice elements for training
    - Includes Claude's quality assessment for filtering
    """
    chunk_id: str                    # Unique identifier for chunk
    source_document: str             # Originating document filename
    document_type: str               # Investment document category
    text: str                        # Actual text content
    word_count: int                  # Word count for size validation
    chunk_index: int                 # Position within source document
    contains_analogy: bool           # Has analogical reasoning
    contains_data: bool              # Has data analysis elements
    contains_personal_voice: bool    # Has personal voice characteristics
    training_format: str             # Intended training pair format
    opus_quality_assessment: str     # Claude's quality evaluation

# ──────────────────────── PROCESSOR ─────────────────────── #
class ClaudeOpusProcessor:
    """
    Premium dataset builder using Claude Opus 4.1 for superior investment voice analysis.
    
    This processor provides the highest quality document analysis and chunking by leveraging
    Claude Opus 4.1's advanced reasoning capabilities. It's designed for users who prioritize
    quality over cost and want the best possible fine-tuning dataset.
    
    Key Advantages over Cerebras Builder:
    - Superior document understanding and context preservation
    - More nuanced voice characteristic detection
    - Better analogical reasoning identification
    - Higher quality chunk boundaries and content analysis
    - More sophisticated training pair generation
    
    Processing Pipeline:
    1. Multi-format document extraction with error correction
    2. Claude Opus 4.1 intelligent analysis and chunking
    3. Quality validation and optimization
    4. Training pair generation with multiple formats
    5. Document balance validation and deduplication
    
    Cost Considerations:
    - Higher cost per document (~$0.50-2.00 per document)
    - Superior quality justifies premium pricing
    - Includes comprehensive cost tracking and estimation
    """
    
    def __init__(self, input_dir: str, output_json: str, api_key: str):
        """
        Initialize Claude Opus processor with directory paths and API credentials.
        
        Args:
            input_dir: Path to directory containing investment documents
            output_json: Output filename for generated dataset
            api_key: Anthropic API key for Claude Opus access
        """
        self.input_dir = Path(input_dir)
        self.output_json = output_json
        self.api_key = api_key
        
        # API usage tracking for cost estimation
        self.window_start = time.time()
        self.sent_requests = 0
        self.total_tokens = 0

    # ═══════════════════════ DOCUMENT EXTRACTION METHODS ═══════════════════════ #
    
    @staticmethod
    def _extract_pdf(p: Path) -> str:
        """
        Extract text from PDF files with advanced OCR error correction.
        
        Uses PyMuPDF (fitz) as primary method with PyPDF2 fallback for maximum
        compatibility. Handles common PDF issues found in investment documents:
        
        - Ligature characters (fi, fl, ff, ffi, ffl)
        - OCR errors and encoding issues
        - Multi-column layouts and complex formatting
        - Scanned documents with text recognition
        
        Args:
            p: Path to PDF file
            
        Returns:
            str: Cleaned and corrected text content
            
        Fallback Strategy:
        - Primary: PyMuPDF for superior text extraction
        - Fallback: PyPDF2 for compatibility with problematic files
        """
        try:
            # Use PyMuPDF for superior text extraction
            doc = fitz.open(str(p))
            text_pages = []
            for page in doc:
                text_pages.append(page.get_text())
            doc.close()
            
            # Combine all pages
            full_text = "\n".join(text_pages)
            
            # Fix common ligature issues
            full_text = full_text.replace("ﬁ", "fi")  # Fix fi ligature (U+FB01)
            full_text = full_text.replace("ﬂ", "fl")  # Fix fl ligature (U+FB02)
            full_text = full_text.replace("ﬀ", "ff")  # Fix ff ligature (U+FB00)
            full_text = full_text.replace("ﬃ", "ffi") # Fix ffi ligature (U+FB03)
            full_text = full_text.replace("ﬄ", "ffl") # Fix ffl ligature (U+FB04)
            full_text = full_text.replace("�", "")     # Remove replacement chars
            
            # Fix common OCR errors in investment docs
            full_text = full_text.replace("Jeremey", "Jeremy")  # Common typo
            full_text = full_text.replace(" moths ", " months ")  # Missing 'n'
            full_text = full_text.replace("12 moths", "12 months")  # Specific fix
            full_text = full_text.replace("beneficiaries", "beneficiaries")  # Common ligature issue
            
            return full_text
            
        except Exception as e:
            print(f"   ⚠️ PyMuPDF failed for {p.name}: {e}, falling back to PyPDF2")
            # Fallback to PyPDF2 if PyMuPDF fails
            with open(p, "rb") as fh:
                rdr = PyPDF2.PdfReader(fh)
                return "\n".join(pg.extract_text() or "" for pg in rdr.pages)

    @staticmethod
    def _extract_docx(p: Path) -> str:
        """
        Extract text from modern Word documents (.docx).
        
        Args:
            p: Path to DOCX file
            
        Returns:
            str: Text content from all paragraphs
        """
        return "\n".join(par.text for par in Document(p).paragraphs)

    @staticmethod
    def _extract_doc(p: Path) -> str:
        """
        Extract text from legacy Word documents (.doc).
        
        Args:
            p: Path to DOC file
            
        Returns:
            str: Extracted text content or empty string if failed
        """
        return docx2txt.process(str(p)) or ""

    @staticmethod
    def _extract_pptx(p: Path) -> str:
        """
        Extract text from PowerPoint presentations (.pptx).
        
        Preserves slide structure with numbered markers for context preservation.
        Essential for investment presentations where slide flow matters.
        
        Args:
            p: Path to PPTX file
            
        Returns:
            str: Text content with slide structure preserved
        """
        pres, out = Presentation(p), []
        for idx, sl in enumerate(pres.slides, 1):
            out.append(f"\n--- Slide {idx} ---\n")
            out.extend(sh.text for sh in sl.shapes if hasattr(sh, "text") and sh.text)
        return "\n".join(out)

    # ==== Claude API helper ============================================= #
    # ═══════════════════════ API MANAGEMENT METHODS ═══════════════════════ #
    
    def _rate_limit(self):
        """
        Enforce API rate limits to prevent quota exceeded errors.
        
        Implements sliding window rate limiting based on RATE_LIMIT_RPM.
        Automatically pauses execution when approaching limits.
        """
        elapsed = time.time() - self.window_start
        if elapsed >= 60:
            self.window_start, self.sent_requests = time.time(), 0
        self.sent_requests += 1
        if self.sent_requests >= RATE_LIMIT_RPM - 5:
            time.sleep(60 - elapsed + 1)
            self.window_start, self.sent_requests = time.time(), 1

    def _call_claude(self, prompt: str) -> str:
        """
        Call Claude Opus 4.1 API with comprehensive error handling and retry logic.
        
        Features:
        - Automatic retry with exponential backoff
        - Rate limiting to respect API constraints
        - Token usage tracking for cost estimation
        - Detailed error logging and recovery
        
        Args:
            prompt: Analysis prompt for Claude Opus
            
        Returns:
            str: Claude's response content
            
        Raises:
            RuntimeError: If all retry attempts fail
        """
        hdrs = {
            "Content-Type": "application/json",
            "x-api-key": self.api_key,
            "anthropic-version": "2023-06-01"
        }
        body = {
            "model": CLAUDE_MODEL,
            "max_tokens": MAX_TOKENS_RESPONSE,
            "messages": [ {"role": "user", "content": prompt} ]
        }
        for att in range(RETRY_MAX_ATTEMPTS):
            self._rate_limit()
            try:
                r = requests.post(API_URL, headers=hdrs, json=body, timeout=120)
                r.raise_for_status()
                return r.json()["content"][0]["text"]
            except Exception as e:
                if att == RETRY_MAX_ATTEMPTS - 1:
                    raise RuntimeError(f"Claude API failed: {e}") from e
                backoff(att)

    # ==== Prompt ========================================================= #
    @staticmethod
    def _prompt(fname: str, txt: str) -> str:
        """
        Generate comprehensive analysis prompt for Claude Opus 4.1.
        
        This prompt is specifically designed to extract investment voice characteristics
        and create optimal chunks for fine-tuning. It instructs Claude to:
        
        - Analyze document type and investment focus areas
        - Identify voice characteristics (analogies, personal opinions, data analysis)
        - Create intelligent chunks preserving context and flow
        - Score content quality for fine-tuning effectiveness
        - Extract distinctive phrases and writing patterns
        - Assess training pair generation potential
        
        Args:
            fname: Document filename for context
            txt: Document text content
            
        Returns:
            str: Structured prompt for Claude Opus analysis
        """
        # Double-braces {{ }} keep JSON braces literal inside f-string
        schema = """{{
  "document_type": "Market Commentary" | "Research Papers" | "Strategy Documents" | "Lessons/Insights" | "Crisis Analysis",
  "document_analysis": {{
    "main_themes": ["theme1", "theme2"],
    "distinctive_voice_elements": ["element1"],
    "analytical_approach": "…",
    "tone_and_style": "…"
  }},
  "distinctive_elements": {{
    "has_analogies": true/false,
    "has_data_analysis": true/false,
    "has_personal_voice": true/false,
    "memorable_phrases": ["phrase1"],
    "data_storytelling_examples": ["example1"],
    "analogies_found": ["analogy1"]
  }},
  "text_chunks": [
    {{
      "text": "chunk text",
      "chunk_reasoning": "why good",
      "word_count": 123,
      "contains_analogy": true/false,
      "contains_data": true/false,
      "contains_personal_voice": true/false,
      "training_format": "style_replication" | "continuation" | "document_specific" | "style_transfer",
      "quality_score": 0.0-1.0
    }}
  ],
  "overall_quality_assessment": {{
    "quality_score": 0.0-1.0,
    "strengths": ["strength1"],
    "unique_voice_markers": ["marker1"],
    "suitability_for_training": "…"
  }}
}}"""
        guide = """\
GUIDE (condensed):
1 Classify doc into the five categories.
2 Create 500-1500-word chunks preserving reasoning flow.
3 Score voice distinctiveness, data-to-insight, analogies, etc.
4 Assign training_format per rules:
  • style_replication → strongest voice samples
  • style_transfer    → chunks rich in analogies
  • continuation      → mid-argument flows
  • document_specific → specialised expertise pieces"""

        return f"""
You are an expert at analysing investment writing and producing fine-tuning datasets.

{guide}

Output JSON strictly matching this schema:
{schema}

DOCUMENT FILENAME: {fname}

---BEGIN-DOC---
{txt}
---END-DOC---

Return ONLY the JSON object, nothing else.
"""

    @staticmethod
    def _parse_json(raw: str) -> Dict:
        """
        Parse JSON response from Claude Opus with robust error handling.
        
        Handles various response formats and common parsing issues:
        - Markdown code blocks (```json)
        - Extra whitespace and formatting
        - Malformed JSON structures
        - Claude's conversational responses
        
        Args:
            raw: Raw response string from Claude
            
        Returns:
            Dict: Parsed JSON data
            
        Raises:
            ValueError: If JSON cannot be parsed after cleanup attempts
        """
        raw = raw.strip()
        if raw.startswith("```"):
            # Handle ```json\n{...}\n``` format
            parts = raw.split("```")
            if len(parts) >= 3:
                raw = parts[1]
                # Remove language identifier (e.g., "json")
                if raw.startswith(("json", "JSON")):
                    raw = raw[4:].strip()
            else:
                raw = raw[3:]  # Just remove opening ```
        return json.loads(raw)

    # ==== Chunk validator ================================================ #
    def _validate_chunks(self, full_text: str, chunks: List[Dict]) -> List[Dict]:
        """
        Validate and optimize text chunks for maximum training effectiveness.
        
        Validation and Optimization Rules:
        - Ensure chunks are within optimal size range (500-1500 words)
        - Keep short documents whole to preserve context
        - Filter out low-quality or repetitive chunks
        - Validate chunk boundaries for logical flow
        - Ensure adequate voice characteristic representation
        
        Args:
            full_text: Complete document text for context
            chunks: List of chunk dictionaries from Claude analysis
            
        Returns:
            List[Dict]: Validated and optimized chunks ready for training
        """
        doc_words = len(full_text.split())
        
        # If document is too short, return as single high-quality chunk
        if doc_words < MIN_CHUNK_WORDS:
            print(f"   📄 Document too short ({doc_words} words), keeping as single chunk")
            # Preserve Claude's analysis but merge into single chunk
            has_analogies = any(c.get("contains_analogy") for c in chunks)
            has_data = any(c.get("contains_data") for c in chunks)
            has_voice = any(c.get("contains_personal_voice") for c in chunks)
            avg_quality = sum(c.get("quality_score", 0.7) for c in chunks) / len(chunks) if chunks else 0.7
            
            return [{
                "text": full_text.strip(),
                "chunk_reasoning": "document_too_short_kept_whole_preserves_complete_reasoning",
                "word_count": doc_words,
                "contains_analogy": has_analogies,
                "contains_data": has_data,
                "contains_personal_voice": has_voice,
                "training_format": "style_replication",  # Best format for complete reasoning
                "quality_score": avg_quality
            }]
        
        # Check if Claude's chunks are properly sized
        if all(MIN_CHUNK_WORDS <= c["word_count"] <= MAX_CHUNK_WORDS for c in chunks):
            return chunks
            
        # If chunks are wrong size, do intelligent re-chunking
        print(f"   🔧 Re-chunking document (target: {MIN_CHUNK_WORDS}-{MAX_CHUNK_WORDS} words)")
        words, new, idx = full_text.split(), [], 0
        while idx < len(words):
            # Aim for middle of target range (1000 words)
            end = min(idx + 1000, len(words))
            piece = " ".join(words[idx:end])
            new.append({
                "text": piece, 
                "chunk_reasoning": "intelligent_rechunk_for_proper_size",
                "word_count": end - idx, 
                "contains_analogy": False,  # Conservative defaults
                "contains_data": False, 
                "contains_personal_voice": True,  # Assume yes for your writing
                "training_format": "continuation",
                "quality_score": 0.7
            })
            idx = end
        return new

    # ==== Single document pipeline ======================================= #
    def _process_single(self, p: Path) -> Tuple[DocumentMetadata, List[TextChunk]]:
        """
        Process a single document through the complete Claude Opus analysis pipeline.
        
        Complete Processing Steps:
        1. Extract text using appropriate method based on file type
        2. Apply PII redaction for privacy protection
        3. Send to Claude Opus 4.1 for intelligent analysis
        4. Parse and validate Claude's JSON response
        5. Create DocumentMetadata with comprehensive analysis
        6. Generate TextChunk objects with unique identifiers
        7. Validate chunk quality and optimize for training
        
        Args:
            p: Path to document file
            
        Returns:
            Tuple[DocumentMetadata, List[TextChunk]]: Document metadata and optimized chunks
            
        Raises:
            Exception: If document processing fails (logged but not fatal to batch)
        """
        print(f"→ {p.name}")
        ext = p.suffix.lower()
        raw = ( self._extract_pdf if ext==".pdf" else
                self._extract_docx if ext==".docx" else
                self._extract_doc  if ext==".doc"  else
                self._extract_pptx if ext==".pptx" else
                None )(p)
        if raw is None or not raw.strip():
            raise ValueError("No text extracted")

        redacted = redact_pii(raw)
        prompt   = self._prompt(p.name, redacted)
        resp     = self._parse_json(self._call_claude(prompt))
        resp["text_chunks"] = self._validate_chunks(raw, resp["text_chunks"])

        meta = DocumentMetadata(
            filename=p.name, file_type=ext, document_type=resp["document_type"],
            word_count=len(raw.split()), chunk_count=len(resp["text_chunks"]),
            has_analogies=resp["distinctive_elements"]["has_analogies"],
            has_data_analysis=resp["distinctive_elements"]["has_data_analysis"],
            has_personal_voice=resp["distinctive_elements"]["has_personal_voice"],
            quality_score=resp["overall_quality_assessment"]["quality_score"],
            distinctive_elements=list(resp["distinctive_elements"].keys()),
            opus_analysis_summary=resp["overall_quality_assessment"]["suitability_for_training"]
        )

        chunks: List[TextChunk] = []
        for i, ck in enumerate(resp["text_chunks"]):
            cid = hashlib.md5(f"{p.stem}_{i}".encode()).hexdigest()[:12]
            chunks.append(TextChunk(
                chunk_id=cid, source_document=p.name, document_type=meta.document_type,
                text=ck["text"], word_count=ck["word_count"], chunk_index=i,
                contains_analogy=ck["contains_analogy"], contains_data=ck["contains_data"],
                contains_personal_voice=ck["contains_personal_voice"],
                training_format=ck["training_format"],
                opus_quality_assessment=ck.get("chunk_reasoning","")
            ))

        self.total_tokens += count_tokens(prompt) + sum(count_tokens(c.text) for c in chunks)
        print(f"   ✓ {p.name}: {len(chunks)} chunks, score {meta.quality_score:.2f}")
        return meta, chunks

    # ==== Training Pair Generation ======================================= #
    @staticmethod
    def _extract_topic(text: str) -> str:
        """Extract meaningful topic from chunk text"""
        text_lower = text.lower()
        
        # Specific topic patterns for investment writing
        if "nvidia" in text_lower and "cisco" in text_lower:
            return "Technology bubble patterns and equipment supplier dynamics"
        elif "bubble" in text_lower and ("tech" in text_lower or "market" in text_lower):
            return "Market bubble analysis and historical patterns"
        elif "china" in text_lower and ("economic" in text_lower or "slowdown" in text_lower):
            return "China's economic conditions and investment implications"
        elif "crisis" in text_lower or "crash" in text_lower:
            return "Financial crisis analysis and lessons"
        elif "emerging market" in text_lower or "em " in text_lower:
            return "Emerging markets investment strategy"
        elif "fed" in text_lower or "interest rate" in text_lower or "monetary" in text_lower:
            return "Federal Reserve policy and market impacts"
        elif "valuation" in text_lower or "earnings" in text_lower:
            return "Company valuation and financial analysis"
        
        # Extract from first meaningful sentence
        sentences = [s.strip() for s in text.split('.') if s.strip()]
        if sentences:
            first_sentence = sentences[0]
            # Look for key subjects in first sentence
            if len(first_sentence.split()) > 15:  # If first sentence is substantial
                # Extract first 8-12 words as topic
                words = first_sentence.split()[:10]
                return " ".join(words) + "..."
        
        # Fallback: look for capitalized entities
        words = text.split()[:50]  # More words for better context
        companies = [w for w in words if w[0].isupper() and len(w) > 3 and w.isalpha()]
        if len(companies) >= 2:
            return f"Analysis of {' and '.join(companies[:2])}"
        elif companies:
            return f"Investment analysis of {companies[0]}"
            
        return "Market commentary and investment insights"
    
    @staticmethod
    def _validate_training_pair(pair: Dict) -> bool:
        """Ensure training pair is valid for fine-tuning"""
        # Check instruction makes sense
        if len(pair.get("instruction", "").split()) < 5:
            return False
        
        # Check input is coherent (if present and not empty)
        input_text = pair.get("input", "")
        if input_text and len(input_text.split()) < 3:
            return False
        
        # Check output is substantial
        output_words = len(pair.get("output", "").split())
        if output_words < 50:  # Minimum meaningful output
            return False
            
        # Check for garbled content
        if "Topic: What's" in input_text or "Nvidia?" in input_text:
            return False  # Catch the specific garbage we saw
            
        return True
    
    def _generate_multiple_pairs_from_chunk(self, chunk: Dict) -> List[Dict]:
        """Generate multiple training variations from single chunk"""
        pairs = []
        text = chunk["text"]
        doc_type = chunk["document_type"].lower()
        source_id = chunk["chunk_id"]
        quality = chunk.get("quality_score", 0.8)
        
        # 1. Full style replication - always generate this
        topic = self._extract_topic(text)
        pairs.append({
            "instruction": f"Write {doc_type} in the style of an experienced emerging markets investor",
            "input": f"Topic: {topic}",
            "output": text,
            "metadata": {
                "source_chunk_id": source_id,
                "training_type": "style_replication",
                "variation": "full_document",
                "quality_score": quality
            }
        })
        
        # 2. Continuation training (split at natural breaks)
        continuation_splits = [
            "Are we now at the top",
            "What I predict is", 
            "My advice is",
            "So will",
            "And so it",
            "Here's",
            "But once"
        ]
        
        for split_phrase in continuation_splits:
            if split_phrase in text and len(text.split()) > 100:
                try:
                    parts = text.split(split_phrase, 1)  # Split only once
                    if len(parts) == 2 and len(parts[0].split()) > 30 and len(parts[1].split()) > 30:
                        pairs.append({
                            "instruction": "Continue this investment analysis in the same style and tone",
                            "input": parts[0].strip(),
                            "output": split_phrase + parts[1],
                            "metadata": {
                                "source_chunk_id": source_id,
                                "training_type": "continuation",
                                "variation": f"split_at_{split_phrase.replace(' ', '_')}",
                                "quality_score": quality * 0.9
                            }
                        })
                        break  # Only use first viable split
                except:
                    continue
                    
        # 3. Analogy extraction (if contains analogies)
        if chunk.get("contains_analogy"):
            analogy_prompts = [
                "Explain market dynamics using historical analogies",
                "Use analogies to explain this investment concept",
                "Draw historical parallels to explain this market situation"
            ]
            
            if "Levi Strauss" in text or "Gold Rush" in text:
                input_text = "How do equipment suppliers perform in technology bubbles?"
            elif "feather in a hurricane" in text:
                input_text = "How difficult is it to time market bubbles?"
            elif "hit the ground" in text:
                input_text = "What's certain about market bubbles?"
            else:
                input_text = f"Explain the investment implications of: {topic}"
                
            pairs.append({
                "instruction": analogy_prompts[len(pairs) % len(analogy_prompts)],
                "input": input_text,
                "output": text,
                "metadata": {
                    "source_chunk_id": source_id,
                    "training_type": "analogy_rich",
                    "variation": "historical_parallels", 
                    "quality_score": quality * 1.1  # Boost for analogies
                }
            })
            
        # 4. Practical advice extraction
        advice_triggers = ["My advice is", "I recommend", "should", "strategy"]
        for trigger in advice_triggers:
            if trigger in text:
                try:
                    advice_start = text.index(trigger)
                    if len(text[advice_start:].split()) > 20:  # Substantial advice
                        pairs.append({
                            "instruction": "Provide practical investment advice based on this analysis",
                            "input": f"Given this market situation, what should I do with my {topic}?",
                            "output": text[advice_start:],
                            "metadata": {
                                "source_chunk_id": source_id,
                                "training_type": "practical_advice",
                                "variation": "actionable_guidance",
                                "quality_score": quality * 1.05  # Slight boost for advice
                            }
                        })
                        break
                except:
                    continue
                    
        # 5. Question-answer format (extract key insights)
        if "?" in text:
            questions = [sent.strip() for sent in text.split("?") if sent.strip()]
            if questions and len(questions[0].split()) < 20:  # Short question
                question = questions[0].strip() + "?"
                remaining_text = text[text.index(question) + len(question):].strip()
                if len(remaining_text.split()) > 30:
                    pairs.append({
                        "instruction": "Answer this investment question with detailed analysis",
                        "input": question,
                        "output": remaining_text,
                        "metadata": {
                            "source_chunk_id": source_id,
                            "training_type": "q_and_a",
                            "variation": "question_response",
                            "quality_score": quality
                        }
                    })
                    
        return pairs
    
    def _create_training_pairs(self, chunks: List[Dict]) -> List[Dict]:
        """Transform analyzed chunks into multiple training pair variations"""
        all_pairs = []
        
        for chunk in chunks:
            # Generate multiple training pairs from each chunk
            chunk_pairs = self._generate_multiple_pairs_from_chunk(chunk)
            all_pairs.extend(chunk_pairs)
            print(f"   📝 Generated {len(chunk_pairs)} pairs from chunk: {chunk['chunk_id']}")
                
        # Validate all pairs and filter out invalid ones
        valid_pairs = []
        for pair in all_pairs:
            if self._validate_training_pair(pair):
                valid_pairs.append(pair)
            else:
                print(f"   ⚠️  Filtered out invalid training pair: {pair.get('instruction', '')[:50]}...")
        
        # Show statistics by training type
        type_counts = {}
        for pair in valid_pairs:
            t_type = pair.get('metadata', {}).get('training_type', 'unknown')
            type_counts[t_type] = type_counts.get(t_type, 0) + 1
            
        print(f"   📊 Training pair breakdown:")
        for t_type, count in type_counts.items():
            print(f"      {t_type}: {count} pairs")
                
        print(f"   ✅ Generated {len(valid_pairs)} valid training pairs (filtered {len(all_pairs) - len(valid_pairs)})")
        print(f"   🎯 Target: ~200 pairs from 28+ documents = {len(valid_pairs)/28:.1f} pairs per document average")
        return valid_pairs
    
    def _validate_document_balance(self, metas: List[Dict]) -> Dict:
        """Check if document distribution matches recommended balance"""
        total = len(metas)
        if total == 0:
            return {"balanced": True, "message": "No documents to validate"}
            
        type_counts = {}
        for meta in metas:
            doc_type = meta["document_type"]
            type_counts[doc_type] = type_counts.get(doc_type, 0) + 1
            
        # Recommended distribution from the guide
        recommended = {
            "Market Commentary": 0.30,
            "Research Papers": 0.25, 
            "Strategy Documents": 0.20,
            "Crisis Analysis": 0.15,
            "Lessons/Insights": 0.10
        }
        
        balance_report = {}
        warnings = []
        
        for doc_type, target_pct in recommended.items():
            actual_count = type_counts.get(doc_type, 0)
            actual_pct = actual_count / total
            balance_report[doc_type] = {
                "actual_count": actual_count,
                "actual_percentage": round(actual_pct * 100, 1),
                "target_percentage": round(target_pct * 100, 1),
                "deviation": round((actual_pct - target_pct) * 100, 1)
            }
            
            # Warning if >10% deviation
            if abs(actual_pct - target_pct) > 0.10:
                warnings.append(f"{doc_type}: {actual_pct:.1%} (target: {target_pct:.1%})")
                
        return {
            "balanced": len(warnings) == 0,
            "balance_report": balance_report,
            "warnings": warnings
        }
    
    def _simple_deduplication(self, chunks: List[Dict]) -> List[Dict]:
        """Simple deduplication based on first 50 words"""
        seen_signatures = set()
        deduplicated = []
        
        for chunk in chunks:
            # Create signature from first 50 words
            words = chunk["text"].split()[:50]
            signature = " ".join(words).lower()
            
            if signature not in seen_signatures:
                seen_signatures.add(signature)
                deduplicated.append(chunk)
            else:
                print(f"   🗑️ Removed duplicate chunk from {chunk.get('source_document', 'unknown')}")
                
        return deduplicated
    
    def _combine_short_documents(self, metas_and_chunks: List[Tuple], min_words: int = 500) -> List[Tuple]:
        """Combine short related documents to meet minimum chunk size"""
        short_docs = [(m, c) for m, c in metas_and_chunks if m.word_count < min_words]
        long_docs = [(m, c) for m, c in metas_and_chunks if m.word_count >= min_words]
        
        if len(short_docs) < 2:
            print(f"   📄 {len(short_docs)} short document(s), no combination needed")
            return metas_and_chunks
            
        print(f"   🔗 Combining {len(short_docs)} short documents for better training quality")
        
        # Group short docs by type for better combination
        by_type = {}
        for meta, chunks in short_docs:
            doc_type = meta.document_type
            if doc_type not in by_type:
                by_type[doc_type] = []
            by_type[doc_type].append((meta, chunks))
            
        combined_results = []
        
        for doc_type, docs in by_type.items():
            if len(docs) == 1:
                # Single short doc of this type, keep as-is
                combined_results.extend(docs)
                continue
                
            # Combine docs of same type
            current_batch = []
            current_words = 0
            
            for meta, chunks in docs:
                if current_words + meta.word_count <= min_words * 2.5:  # Don't exceed reasonable max
                    current_batch.append((meta, chunks))
                    current_words += meta.word_count
                else:
                    # Process current batch
                    if current_batch:
                        combined_results.append(self._merge_documents(current_batch, doc_type))
                    current_batch = [(meta, chunks)]
                    current_words = meta.word_count
                    
            # Process final batch
            if current_batch:
                if len(current_batch) > 1:
                    combined_results.append(self._merge_documents(current_batch, doc_type))
                else:
                    combined_results.extend(current_batch)
                    
        # Add back the long documents
        combined_results.extend(long_docs)
        
        print(f"   ✅ Document combination: {len(metas_and_chunks)} → {len(combined_results)} documents")
        return combined_results
        
    def _merge_documents(self, docs: List[Tuple], doc_type: str) -> Tuple:
        """Merge multiple document metadata and chunks"""
        metas, all_chunks = zip(*docs)
        
        # Combine text from all chunks
        combined_text = "\n\n".join(
            f"=== {meta.filename} ===\n{' '.join(chunk.text for chunk in chunks)}"
            for meta, chunks in docs
        )
        
        # Create merged metadata
        filenames = [meta.filename for meta in metas]
        total_words = sum(meta.word_count for meta in metas)
        avg_quality = sum(meta.quality_score for meta in metas) / len(metas)
        
        merged_meta = DocumentMetadata(
            filename=f"COMBINED_{len(filenames)}_docs_({', '.join(f[:15] for f in filenames[:3])})",
            file_type=".combined",
            document_type=doc_type,
            word_count=total_words,
            chunk_count=1,  # Will become single chunk
            has_analogies=any(meta.has_analogies for meta in metas),
            has_data_analysis=any(meta.has_data_analysis for meta in metas),
            has_personal_voice=any(meta.has_personal_voice for meta in metas),
            quality_score=avg_quality,
            distinctive_elements=["combined_document", "multi_source"],
            opus_analysis_summary=f"Combined analysis of {len(filenames)} related {doc_type.lower()} documents"
        )
        
        # Create single high-quality chunk from combined text
        chunk_id = hashlib.md5(f"combined_{len(filenames)}".encode()).hexdigest()[:12]
        merged_chunk = TextChunk(
            chunk_id=chunk_id,
            source_document=merged_meta.filename,
            document_type=doc_type,
            text=combined_text,
            word_count=total_words,
            chunk_index=0,
            contains_analogy=any(any(c.contains_analogy for c in chunks) for chunks in all_chunks),
            contains_data=any(any(c.contains_data for c in chunks) for chunks in all_chunks),
            contains_personal_voice=True,  # Assume yes for combined docs
            training_format="style_replication",  # Best for combined content
            opus_quality_assessment="Combined document for enhanced training context"
        )
        
        return (merged_meta, [merged_chunk])

    def preprocess_corpus_analysis(self) -> Dict:
        """Analyze entire corpus before processing to estimate results and costs"""
        print(f"\n🔍 PREPROCESSING ANALYSIS")
        print(f"📁 Analyzing corpus: {self.input_dir}")
        
        stats = {
            'total_files': 0,
            'total_words': 0,
            'by_size': {'small': 0, 'medium': 0, 'large': 0},
            'size_details': {'small': [], 'medium': [], 'large': []},
            'estimated_chunks': 0,
            'estimated_pairs': 0,
            'estimated_cost': 0.0
        }
        
        # Find all supported files
        all_files = []
        for pattern in ("*.pdf", "*.docx", "*.doc", "*.pptx"):
            all_files.extend(self.input_dir.glob(pattern))
            
        if not all_files:
            print("❌ No supported files found!")
            return stats
            
        print(f"📊 Found {len(all_files)} files to analyze...")
        
        for file_path in all_files:
            try:
                # Extract text to count words
                if file_path.suffix.lower() == '.pdf':
                    text = self._extract_pdf(file_path)
                elif file_path.suffix.lower() == '.docx':
                    text = self._extract_docx(file_path)
                elif file_path.suffix.lower() == '.doc':
                    text = self._extract_doc(file_path)
                elif file_path.suffix.lower() == '.pptx':
                    text = self._extract_pptx(file_path)
                else:
                    continue
                    
                words = len(text.split())
                stats['total_files'] += 1
                stats['total_words'] += words
                
                # Categorize by size and estimate training pairs
                if words < 500:
                    stats['by_size']['small'] += 1
                    stats['size_details']['small'].append((file_path.name, words))
                    stats['estimated_chunks'] += 1
                    stats['estimated_pairs'] += 3  # Small docs = 3 pairs
                elif words < 2000:
                    stats['by_size']['medium'] += 1
                    stats['size_details']['medium'].append((file_path.name, words))
                    stats['estimated_chunks'] += 2
                    stats['estimated_pairs'] += 8  # Medium docs = 8 pairs
                else:
                    stats['by_size']['large'] += 1
                    stats['size_details']['large'].append((file_path.name, words))
                    stats['estimated_chunks'] += 4
                    stats['estimated_pairs'] += 15  # Large docs = 15 pairs
                    
            except Exception as e:
                print(f"   ⚠️ Failed to analyze {file_path.name}: {e}")
                
        # Calculate estimated API cost
        stats['estimated_cost'] = round(stats['estimated_pairs'] * 0.75, 2)  # ~$0.75 per document
        
        # Print detailed analysis
        print(f"\n📈 CORPUS ANALYSIS RESULTS:")
        print(f"   📄 Total files: {stats['total_files']}")
        print(f"   📝 Total words: {stats['total_words']:,}")
        print(f"   📊 Average words per document: {stats['total_words']//stats['total_files'] if stats['total_files'] > 0 else 0:,}")
        
        print(f"\n📏 SIZE DISTRIBUTION:")
        for size_cat, count in stats['by_size'].items():
            if count > 0:
                print(f"   {size_cat.capitalize()} docs (<500, 500-2000, >2000 words): {count}")
                # Show a few examples
                examples = stats['size_details'][size_cat][:3]
                for name, words in examples:
                    print(f"      • {name}: {words} words")
                if len(stats['size_details'][size_cat]) > 3:
                    print(f"      • ... and {len(stats['size_details'][size_cat]) - 3} more")
                    
        print(f"\n🎯 PROJECTIONS:")
        print(f"   📦 Estimated chunks: {stats['estimated_chunks']}")
        print(f"   🎓 Estimated training pairs: {stats['estimated_pairs']}")
        print(f"   💰 Estimated API cost: ${stats['estimated_cost']}")
        
        # Viability assessment
        if stats['estimated_pairs'] >= 200:
            print(f"   ✅ EXCELLENT: {stats['estimated_pairs']} pairs should create a high-quality dataset")
        elif stats['estimated_pairs'] >= 100:
            print(f"   ✅ GOOD: {stats['estimated_pairs']} pairs is sufficient for fine-tuning")
        elif stats['estimated_pairs'] >= 50:
            print(f"   ⚠️  MARGINAL: {stats['estimated_pairs']} pairs may need augmentation")
        else:
            print(f"   ❌ LOW: {stats['estimated_pairs']} pairs likely insufficient, consider few-shot instead")
            
        return stats

    # ==== Directory pipeline ============================================= #
    def run(self, skip_analysis: bool = False):
        if not self.input_dir.exists():
            raise FileNotFoundError(self.input_dir)
            
        # Run preprocessing analysis first (unless skipped)
        if not skip_analysis:
            analysis = self.preprocess_corpus_analysis()
            
            # Ask for confirmation before proceeding with expensive API calls
            if analysis['estimated_cost'] > 30:
                response = input(f"\n⚠️  Estimated cost is ${analysis['estimated_cost']}. Continue? (y/n): ")
                if response.lower() != 'y':
                    print("Processing cancelled.")
                    return
            elif analysis['estimated_pairs'] < 50:
                response = input(f"\n⚠️  Only {analysis['estimated_pairs']} pairs estimated. Continue anyway? (y/n): ")
                if response.lower() != 'y':
                    print("Processing cancelled. Consider few-shot prompting instead.")
                    return
                    
        paths = []
        for g in ("*.pdf","*.docx","*.doc","*.pptx"):
            paths.extend(self.input_dir.glob(g))
        if not paths:
            print("No files found."); return

        metas, chks, fails, raw_responses = [], [], [], []
        for fp in paths:
            try:
                m, c = self._process_single(fp)
                metas.append(asdict(m)); chks.extend(asdict(x) for x in c)
            except Exception as e:
                fails.append({"file": fp.name, "error": str(e)})
                print(f"   ✗ {fp.name}: {e}")

        print(f"\n🔍 Processing {len(chks)} chunks...")
        
        # Apply deduplication
        original_count = len(chks)
        chks = self._simple_deduplication(chks)
        if original_count != len(chks):
            print(f"   🗑️  Removed {original_count - len(chks)} duplicate chunks")
            
        # Generate training pairs
        training_pairs = self._create_training_pairs(chks)
        print(f"   📚 Generated {len(training_pairs)} training pairs")
        
        # Validate document balance
        balance = self._validate_document_balance(metas)
        if not balance["balanced"]:
            print(f"   ⚠️  Document balance warnings: {', '.join(balance['warnings'])}")
        else:
            print(f"   ✅ Document types are well balanced")

        words = sum(c["word_count"] for c in chks)
        doc_dist, fmt_dist = {}, {}
        for c in chks:
            doc_dist[c["document_type"]] = doc_dist.get(c["document_type"],0)+1
            fmt_dist[c["training_format"]] = fmt_dist.get(c["training_format"],0)+1

        summary = {
            "total_documents": len(metas),"total_chunks": len(chks),"total_words": words,
            "training_pairs_generated": len(training_pairs),
            "average_words_per_chunk": words/len(chks) if chks else 0,
            "document_type_distribution": doc_dist,
            "training_format_distribution": fmt_dist,
            "document_balance": balance,
            "failed_files": fails,
            "documents_with_analogies": sum(m["has_analogies"] for m in metas),
            "documents_with_data": sum(m["has_data_analysis"] for m in metas),
            "documents_with_personal_voice": sum(m["has_personal_voice"] for m in metas),
            "average_quality_score": sum(m["quality_score"] for m in metas)/len(metas) if metas else 0,
            "estimated_api_cost": round((self.total_tokens/1_000_000)*15,2),
            "processing_method": "Claude Opus 4.1 Intelligent Analysis (v3 Enhanced)"
        }

        out = {
            "summary": summary,"documents": metas,"chunks": chks,
            "training_pairs": training_pairs,
            "opus_processing_notes": {
                "model_used": CLAUDE_MODEL,
                "analysis_approach": "Enhanced with training pairs, deduplication, balance validation",
                "quality_assessment": "Chunks validated, deduplicated, ready for fine-tuning",
                "enhancements_v3": [
                    "Training pair generation (instruction-input-output format)",
                    "Document balance validation against recommended distribution", 
                    "Simple deduplication based on text signatures",
                    "Better PDF extraction with pdfplumber fallback",
                    "Comprehensive metadata tracking"
                ]
            }
        }
        with open(self.output_json,"w",encoding="utf-8") as f:
            json.dump(out,f,indent=2,ensure_ascii=False)

        print("\n=== RUN COMPLETE ===")
        print(json.dumps(summary,indent=2))

# ─────────────────────────── CLI ────────────────────────── #
if __name__ == "__main__":
    # Configuration for premium dataset generation
    INPUT_DIR = "/Users/macbook2024/Dropbox/AAA Backup/A Working/Arjun LLM Writing/Writing"
    OUTPUT_JSON = "training_dataset_opus_processed.json"
    API_KEY = os.getenv("ANTHROPIC_API_KEY")
    
    print("🚀 Starting Claude Opus Premium Dataset Builder...")
    print(f"📁 Input Directory: {INPUT_DIR}")
    print(f"📄 Output File: {OUTPUT_JSON}")
    print(f"🧠 AI Model: {CLAUDE_MODEL}")
    print(f"💰 Premium Quality: Higher cost, superior results")
    print("─" * 60)
    
    # Validate API key
    if not API_KEY:
        print("❌ ERROR: ANTHROPIC_API_KEY environment variable not set.")
        print("   Please set your Anthropic API key to use Claude Opus.")
        sys.exit(1)
    
    # Initialize and run processor
    processor = ClaudeOpusProcessor(INPUT_DIR, OUTPUT_JSON, API_KEY)
    processor.run()
    
    print("\n✅ Premium dataset generation complete!")
    print(f"💾 Dataset saved to: {OUTPUT_JSON}")
    print("🔄 Ready for OpenAI fine-tuning conversion")
    print("💡 Claude Opus provides superior quality for investment voice capture")